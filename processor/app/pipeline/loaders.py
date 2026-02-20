from __future__ import annotations

from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation


def extract_pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        body = page.extract_text() or ""
        pages.append(f"Page {i}\n{body}")
    return "\n".join(pages)


def extract_pptx_text(path: Path) -> str:
    prs = Presentation(str(path))
    slides_text: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = [f"Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text.strip())
        slides_text.append("\n".join(lines))
    return "\n\n".join(slides_text)
